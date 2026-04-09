"""
Patch script: adds a Database Schema slide to SISM_Final_Presentation.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import copy, sys

DARK_BG      = RGBColor(0x0D, 0x14, 0x2B)
CARD_BG      = RGBColor(0x13, 0x22, 0x3F)
ACCENT_BLUE  = RGBColor(0x3B, 0x82, 0xF6)
ACCENT_CYAN  = RGBColor(0x06, 0xB6, 0xD4)
ACCENT_GREEN = RGBColor(0x10, 0xB9, 0x81)
ACCENT_AMBER = RGBColor(0xF5, 0x9E, 0x0B)
ACCENT_ROSE  = RGBColor(0xF4, 0x3F, 0x5E)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY   = RGBColor(0x94, 0xA3, 0xB8)
PANEL_BORDER = RGBColor(0x1E, 0x3A, 0x5F)

PPTX_PATH = "SISM_Final_Presentation.pptx"
prs = Presentation(PPTX_PATH)
blank = prs.slide_layouts[6]


def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def rect(slide, l, t, w, h, fill=None, line=None, lw=1):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.fill.solid() if fill else s.fill.background()
    if fill: s.fill.fore_color.rgb = fill
    if line: s.line.color.rgb = line; s.line.width = Pt(lw)
    else: s.line.fill.background()
    return s


def rrect(slide, l, t, w, h, fill=None, line=None, lw=1):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.adjustments[0] = 0.04
    s.fill.solid() if fill else s.fill.background()
    if fill: s.fill.fore_color.rgb = fill
    if line: s.line.color.rgb = line; s.line.width = Pt(lw)
    else: s.line.fill.background()
    return s


def tb(slide, text, l, t, w, h, sz=12, bold=False, color=WHITE,
       align=PP_ALIGN.LEFT, italic=False):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run()
    run.text = text; run.font.size = Pt(sz)
    run.font.bold = bold; run.font.italic = italic
    run.font.color.rgb = color; run.font.name = "Segoe UI"
    return box


# ─── Build the new slide ──────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank)
set_bg(slide, DARK_BG)
rect(slide, Inches(0), Inches(0), Inches(13.33), Inches(0.08), fill=ACCENT_BLUE)

# Header
rect(slide, Inches(0.5), Inches(0.45), Inches(12.33), Inches(0.04), fill=ACCENT_BLUE)
tb(slide, "Database Schema", Inches(0.5), Inches(0.55), Inches(10), Inches(0.65),
   sz=32, bold=True, color=WHITE)
tb(slide, "9 SQLAlchemy ORM Models → 9 SQLite Tables  |  SQLite file: backend/sism.db",
   Inches(0.5), Inches(1.15), Inches(12.5), Inches(0.38), sz=13, color=LIGHT_GRAY)

# ── Table definitions (name, accent_color, [(col_name, type, note)])
tables = [
    ("perishable_items", ACCENT_CYAN, [
        ("id", "INT PK", "Auto-increment primary key"),
        ("sku", "VARCHAR", "Unique product code e.g. 'AVO-001'"),
        ("name / category", "VARCHAR", "e.g. 'Avocado' / 'tropical_fruit'"),
        ("shelf_life_days", "INT", "Max usable days from received date"),
        ("ripeness_curve", "VARCHAR", "fast_sigmoid / slow_sigmoid / linear / flat"),
        ("ripeness_peak_day", "INT", "Day number when fruit hits peak ripeness"),
        ("zone", "VARCHAR", "ambient / refrigerated / cold_chain / freezer"),
        ("reorder_point_kg", "FLOAT", "Stock threshold to trigger procurement"),
    ]),
    ("storage_containers", ACCENT_GREEN, [
        ("id", "INT PK", "Auto-increment primary key"),
        ("container_id", "VARCHAR", "e.g. 'REF-A', 'AMB-B'"),
        ("zone_type", "VARCHAR", "ambient / refrigerated / cold_chain / freezer"),
        ("capacity_kg / current_load_kg", "FLOAT", "Max capacity and live fill level"),
        ("rows / cols / depths", "INT", "3D physical grid dimensions"),
        ("temp_c", "FLOAT", "IoT-measured temperature in Celsius"),
    ]),
    ("inventory_batches", ACCENT_ROSE, [
        ("id / batch_id", "INT/VARCHAR", "Primary key + human-readable ID e.g. BAT-A1F3"),
        ("item_id", "INT FK", "→ perishable_items.id"),
        ("container_id", "INT FK", "→ storage_containers.id"),
        ("quantity_kg", "FLOAT", "Weight of batch in kilograms"),
        ("received_date", "DATE", "Date batch arrived at warehouse"),
        ("expected_ripeness_date", "DATE", "Calculated by Ripeness Predictor engine"),
        ("expiry_date", "DATE", "Hard expiry limit (never use after this)"),
        ("ripeness_score", "FLOAT", "0.0=unripe · 1.0=peak · 2.0+=overripe"),
        ("status", "VARCHAR", "in_stock / reserved / expired / wasted / at_risk"),
        ("row / col / depth", "INT", "Physical 3D coordinates in container grid"),
    ]),
    ("delivery_orders", ACCENT_AMBER, [
        ("id / order_id", "INT/VARCHAR", "Primary key + human ID e.g. ORD-B2C9"),
        ("client_name / city", "VARCHAR", "Customer and delivery destination"),
        ("delivery_date", "DATE", "Scheduled delivery date"),
        ("status", "VARCHAR", "pending / confirmed / fulfilled / cancelled"),
        ("priority", "VARCHAR", "low / normal / high / urgent"),
    ]),
    ("order_items", ACCENT_BLUE, [
        ("order_id", "INT FK", "→ delivery_orders.id"),
        ("item_id", "INT FK", "→ perishable_items.id"),
        ("quantity_kg", "FLOAT", "Amount ordered by client"),
        ("fulfilled_kg", "FLOAT", "Amount actually delivered (0 if pending)"),
        ("unit_price", "FLOAT", "Price per kg at time of order"),
    ]),
    ("procurement_recommendations", ACCENT_CYAN, [
        ("rec_id", "VARCHAR", "Unique recommendation ID"),
        ("item_id", "INT FK", "→ perishable_items.id"),
        ("recommended_qty_kg", "FLOAT", "AI-calculated order quantity in kg"),
        ("order_by_date", "DATE", "Deadline to place supplier order"),
        ("priority", "VARCHAR", "normal / urgent"),
        ("ai_generated", "BOOL", "True = generated by Gemini, False = manual"),
        ("status", "VARCHAR", "pending / approved / rejected / ordered"),
    ]),
    ("ai_insights", ACCENT_AMBER, [
        ("insight_id", "VARCHAR", "Unique ID e.g. INS-A3F9"),
        ("insight_type", "VARCHAR", "expiry_alert / stock_low / placement_conflict / demand_gap"),
        ("title / message", "VARCHAR/TEXT", "Short title + full description of the alert"),
        ("severity", "VARCHAR", "info / warning / critical"),
        ("resolved", "BOOL", "True = dismissed by manager"),
    ]),
    ("storage_placements", ACCENT_GREEN, [
        ("batch_id", "INT FK", "→ inventory_batches.id"),
        ("container_id", "INT FK", "→ storage_containers.id"),
        ("row / col / depth", "INT", "Physical coordinates in 3D grid"),
        ("conflict_flag", "BOOL", "True = LIFO conflict detected at this position"),
        ("conflict_reason", "TEXT", "Explanation of why conflict exists"),
    ]),
    ("audit_logs", ACCENT_ROSE, [
        ("entity_type / entity_id", "VARCHAR", "Which table and which record was changed"),
        ("action", "VARCHAR", "CREATE / UPDATE / DELETE / OPTIMIZE / AI_RECOMMEND"),
        ("old_value / new_value", "JSON", "Before and after state snapshot"),
        ("performed_by", "VARCHAR", "'system' or manager username"),
        ("timestamp", "DATETIME", "Exact time the change occurred"),
    ]),
]

# Layout: 3 columns of 3 tables
col_positions = [Inches(0.25), Inches(4.55), Inches(8.85)]
row_positions = [Inches(1.65), Inches(3.55), Inches(5.45)]
table_w = Inches(4.2)
table_h = Inches(1.75)

for idx, (tname, col, cols) in enumerate(tables):
    lft = col_positions[idx % 3]
    top = row_positions[idx // 3]

    # card
    rrect(slide, lft, top, table_w, table_h, fill=CARD_BG, line=col, lw=1.5)
    rect(slide, lft, top, table_w, Inches(0.07), fill=col)

    # table name
    tb(slide, tname, lft + Inches(0.12), top + Inches(0.1),
       table_w - Inches(0.2), Inches(0.35), sz=11, bold=True, color=col)

    # columns
    for ci, (cname, dtype, note) in enumerate(cols[:4]):
        cy = top + Inches(0.5) + ci * Inches(0.3)
        tb(slide, cname, lft + Inches(0.12), cy, Inches(1.55), Inches(0.28),
           sz=8.5, bold=True, color=WHITE)
        tb(slide, dtype, lft + Inches(1.68), cy, Inches(1.0), Inches(0.28),
           sz=8, color=ACCENT_CYAN)
        tb(slide, note, lft + Inches(2.7), cy, Inches(1.4), Inches(0.28),
           sz=7.5, color=LIGHT_GRAY, italic=True)

    if len(cols) > 4:
        tb(slide, f"+ {len(cols)-4} more columns...",
           lft + Inches(0.12), top + Inches(1.55), table_w, Inches(0.25),
           sz=8, color=LIGHT_GRAY, italic=True)


# ── Move slide to position 2 (after title slide) ─────────────────────────────
xml_slides = prs.slides._sldIdLst
# The new slide is at the end; move it to index 1
new_slide_elem = xml_slides[-1]
xml_slides.remove(new_slide_elem)
xml_slides.insert(1, new_slide_elem)

prs.save(PPTX_PATH)
print(f"\n✅  Database Schema slide inserted at position 2 in {PPTX_PATH}")
print(f"   Total slides now: {len(prs.slides)}\n")
